import os
import io
import time
import requests
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from PIL import Image, ImageTk
import cv2
import numpy as np
import threading
import tempfile
import subprocess
import platform
import psutil
import json
import base64
import logging
import sys
import datetime
import traceback
from io import BytesIO

class PowerPointTranslator:
    @staticmethod
    def setup_logging():
        """로깅 시스템 설정"""
        log_file = f"translator_debug_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.log"
        
        # 루트 로거 설정
        root_logger = logging.getLogger()
        root_logger.setLevel(logging.DEBUG)
        
        # 파일 핸들러 설정
        file_handler = logging.FileHandler(log_file, encoding='utf-8')
        file_handler.setLevel(logging.DEBUG)
        
        # 콘솔 핸들러 설정
        console_handler = logging.StreamHandler(sys.stdout)
        console_handler.setLevel(logging.INFO)
        
        # 포맷 설정
        formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
        file_handler.setFormatter(formatter)
        console_handler.setFormatter(formatter)
        
        # 핸들러 추가
        root_logger.addHandler(file_handler)
        root_logger.addHandler(console_handler)
        
        return log_file
    
    def __init__(self, root):
        # 로깅 설정
        self.logger = logging.getLogger(__name__)
        self.log_file = self.setup_logging()
        self.logger.info("프로그램 시작")
        
        self.root = root
        self.root.title("Powerpoint Image Translator")
        self.root.geometry("850x900")
        
        # 모델 초기화 상태 추적
        self.models_initialized = False
        
        # 상단 프레임 (로고와 제목을 위한 프레임)
        self.top_frame = tk.Frame(root)
        self.top_frame.grid(row=0, column=0, columnspan=2, padx=20, pady=20, sticky="ew")
        
        # 그리드 설정
        self.top_frame.columnconfigure(0, weight=1)  # 로고 열
        self.top_frame.columnconfigure(1, weight=2)  # 제목 열
        self.top_frame.columnconfigure(2, weight=1)  # 빈 공간 열
        
        # LINE studio 로고 로드
        try:
            self.logo_photo = tk.PhotoImage(file="line-studio-logo.png")
            self.logo_label = tk.Label(self.top_frame, image=self.logo_photo)
            self.logo_label.grid(row=0, column=0, sticky="w")
        except Exception as e:
            self.logger.warning(f"로고 로드 오류: {e}")
            # 로고 로드 실패 시 텍스트로 대체
            self.logo_label = tk.Label(self.top_frame, text="LINE studio", fg="#8CC63F", font=("Arial", 24, "bold"))
            self.logo_label.grid(row=0, column=0, sticky="w")
        
        # 제목 (가운데 정렬)
        self.title_label = tk.Label(self.top_frame, text="Powerpoint Image Translator", font=("Arial", 24))
        self.title_label.grid(row=0, column=1)
        
        # 파일 프레임
        self.file_frame = tk.Frame(root, highlightbackground="#ddd", highlightthickness=1, padx=10, pady=10)
        self.file_frame.grid(row=1, column=0, columnspan=2, padx=50, pady=20, sticky="ew")
        
        # 파일 경로 레이블
        self.file_path_label = tk.Label(self.file_frame, text="파일 경로:")
        self.file_path_label.grid(row=0, column=0, padx=(0, 10), pady=10, sticky="w")
        
        # 파일 경로 입력 상자
        self.file_path_var = tk.StringVar()
        self.file_path_entry = tk.Entry(self.file_frame, textvariable=self.file_path_var, width=60)
        self.file_path_entry.grid(row=0, column=1, padx=5, pady=10, sticky="ew")
        
        # 찾아보기 버튼
        self.browse_button = tk.Button(self.file_frame, text="찾아보기", 
                                      command=self.select_file)
        self.browse_button.grid(row=0, column=2, padx=5, pady=10)
        
        # 그리드 컬럼 설정
        self.file_frame.columnconfigure(1, weight=1)
        
        # 정보 및 진행상황 프레임 (2개의 LabelFrame 가로 배치)
        self.info_progress_frame = tk.Frame(root)
        self.info_progress_frame.grid(row=2, column=0, columnspan=2, padx=50, pady=10, sticky="ew")
        self.info_progress_frame.columnconfigure(0, weight=1)
        self.info_progress_frame.columnconfigure(1, weight=1)
        
        # 파일 정보 표시 영역
        self.info_frame = tk.LabelFrame(self.info_progress_frame, text="파일 정보", padx=10, pady=10)
        self.info_frame.grid(row=0, column=0, padx=(0, 5), pady=5, sticky="nsew")
        
        # 파일 정보 라벨
        self.file_name_label = tk.Label(self.info_frame, text="파일 이름: 선택된 파일 없음")
        self.file_name_label.grid(row=0, column=0, sticky="w", pady=3)
        
        self.slide_count_label = tk.Label(self.info_frame, text="슬라이드 수: -")
        self.slide_count_label.grid(row=1, column=0, sticky="w", pady=3)
        
        self.text_count_label = tk.Label(self.info_frame, text="텍스트 요소 수: -")
        self.text_count_label.grid(row=2, column=0, sticky="w", pady=3)
        
        self.image_count_label = tk.Label(self.info_frame, text="이미지 요소 수: -")
        self.image_count_label.grid(row=3, column=0, sticky="w", pady=3)
        
        self.total_elements_label = tk.Label(self.info_frame, text="총 번역 요소: -")
        self.total_elements_label.grid(row=4, column=0, sticky="w", pady=3)
        
        # 진행상황 표시 영역
        self.progress_frame = tk.LabelFrame(self.info_progress_frame, text="진행 상황", padx=10, pady=10)
        self.progress_frame.grid(row=0, column=1, padx=(5, 0), pady=5, sticky="nsew")
        
        # 진행상황 라벨
        self.current_slide_label = tk.Label(self.progress_frame, text="현재 슬라이드: -")
        self.current_slide_label.grid(row=0, column=0, sticky="w", pady=3)
        
        self.current_task_label = tk.Label(self.progress_frame, text="현재 작업: -")
        self.current_task_label.grid(row=1, column=0, sticky="w", pady=3)
        
        self.translated_items_label = tk.Label(self.progress_frame, text="번역된 요소: -")
        self.translated_items_label.grid(row=2, column=0, sticky="w", pady=3)
        
        self.remaining_items_label = tk.Label(self.progress_frame, text="남은 요소: -")
        self.remaining_items_label.grid(row=3, column=0, sticky="w", pady=3)
        
        # Ollama 상태 표시
        self.ollama_status_frame = tk.LabelFrame(root, text="Ollama 상태", padx=10, pady=10)
        self.ollama_status_frame.grid(row=3, column=0, columnspan=2, padx=50, pady=10, sticky="ew")
        
        # Ollama 설치 상태
        self.ollama_installed_label = tk.Label(self.ollama_status_frame, text="설치 상태: 확인 중...")
        self.ollama_installed_label.grid(row=0, column=0, sticky="w", pady=2)
        
        # Ollama 실행 상태
        self.ollama_running_label = tk.Label(self.ollama_status_frame, text="실행 상태: 확인 중...")
        self.ollama_running_label.grid(row=1, column=0, sticky="w", pady=2)
        
        # Ollama 포트
        self.ollama_port_label = tk.Label(self.ollama_status_frame, text="Ollama 포트: 확인 중...")
        self.ollama_port_label.grid(row=2, column=0, sticky="w", pady=2)
        
        # Ollama 상태 확인 버튼
        self.check_ollama_button = tk.Button(self.ollama_status_frame, text="상태 확인", 
                                          command=self.check_ollama_status)
        self.check_ollama_button.grid(row=0, column=1, rowspan=3, padx=20)
        
        # 번역 옵션 프레임
        self.options_frame = tk.LabelFrame(root, text="번역 옵션", padx=10, pady=10)
        self.options_frame.grid(row=4, column=0, columnspan=2, padx=50, pady=10, sticky="ew")
        
        # 언어 목록 정의
        self.languages = ["한국어", "일본어", "영어", "중국어번체", "중국어간체", "태국어", "스페인어", "프랑스어"]
        
        # 원본 언어 선택
        self.source_label = tk.Label(self.options_frame, text="원본 언어:")
        self.source_label.grid(row=0, column=0, padx=10, pady=10)
        
        self.source_lang = tk.StringVar(value="일본어")
        self.source_combo = ttk.Combobox(self.options_frame, textvariable=self.source_lang, 
                                        values=self.languages)
        self.source_combo.grid(row=0, column=1, padx=10, pady=10)
        
        # 화살표 레이블
        self.arrow_label = tk.Label(self.options_frame, text="⟷")
        self.arrow_label.grid(row=0, column=2, padx=10, pady=10)
        
        # 번역 언어 선택
        self.target_label = tk.Label(self.options_frame, text="번역 언어:")
        self.target_label.grid(row=0, column=3, padx=10, pady=10)
        
        self.target_lang = tk.StringVar(value="한국어")
        self.target_combo = ttk.Combobox(self.options_frame, textvariable=self.target_lang, 
                                        values=self.languages)
        self.target_combo.grid(row=0, column=4, padx=10, pady=10)
        
        # Vision 모델 선택
        self.vision_model_label = tk.Label(self.options_frame, text="Vision 모델:")
        self.vision_model_label.grid(row=1, column=0, padx=10, pady=10)
        
        self.vision_model_var = tk.StringVar(value="")
        self.vision_model_combo = ttk.Combobox(self.options_frame, textvariable=self.vision_model_var, state="readonly")
        self.vision_model_combo.grid(row=1, column=1, padx=10, pady=10, sticky="ew")
        
        # Text 모델 선택
        self.text_model_label = tk.Label(self.options_frame, text="Text 모델:")
        self.text_model_label.grid(row=1, column=3, padx=10, pady=10)
        
        self.text_model_var = tk.StringVar(value="gemma3:12b")
        self.text_model_combo = ttk.Combobox(self.options_frame, textvariable=self.text_model_var, state="readonly")
        self.text_model_combo.grid(row=1, column=4, padx=10, pady=10, sticky="ew")
        
        # URL 입력
        self.url_label = tk.Label(self.options_frame, text="Ollama URL:")
        self.url_label.grid(row=2, column=0, padx=10, pady=10)
        
        self.url_var = tk.StringVar(value="http://localhost:11434")
        self.url_entry = tk.Entry(self.options_frame, textvariable=self.url_var, width=40)
        self.url_entry.grid(row=2, column=1, columnspan=4, padx=10, pady=10, sticky="ew")
        
        # 번역 시작/중지 버튼
        self.buttons_frame = tk.Frame(root)
        self.buttons_frame.grid(row=5, column=0, columnspan=2, padx=50, pady=10)
        
        self.start_button = tk.Button(self.buttons_frame, text="번역 시작", 
                                    bg="#4999E9", fg="white", width=20, height=2,
                                    command=self.start_translation)
        self.start_button.grid(row=0, column=0, padx=10)
        
        self.stop_button = tk.Button(self.buttons_frame, text="번역 중지", 
                                    bg="#CCCCCC", fg="white", width=20, height=2,
                                    command=self.stop_translation,
                                    state=tk.DISABLED)
        self.stop_button.grid(row=0, column=1, padx=10)
        
        # 진행 상황 프레임
        self.progress_bar_frame = tk.Frame(root)
        self.progress_bar_frame.grid(row=6, column=0, columnspan=2, padx=50, pady=5, sticky="ew")
        
        # 진행 상황 막대
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(self.progress_bar_frame, variable=self.progress_var, maximum=100, length=600)
        self.progress_bar.grid(row=0, column=0, sticky="ew")
        
        # 진행율과 시간 표시 레이블
        self.progress_label = tk.Label(self.progress_bar_frame, text="0% (0:00 / 계산 중)")
        self.progress_label.grid(row=0, column=1, padx=10)
        
        # 상태 메시지
        self.status_label = tk.Label(root, text="준비 완료")
        self.status_label.grid(row=7, column=0, columnspan=2, padx=50, pady=5, sticky="w")
        
        # 로그 프레임
        self.log_frame = tk.LabelFrame(root, text="로그", padx=10, pady=10)
        self.log_frame.grid(row=8, column=0, columnspan=2, padx=50, pady=10, sticky="ew")
        
        # 로그 텍스트 영역
        self.log_text = tk.Text(self.log_frame, height=10, width=80, wrap=tk.WORD)
        self.log_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        # 스크롤바
        log_scrollbar = tk.Scrollbar(self.log_frame, command=self.log_text.yview)
        log_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.log_text.config(yscrollcommand=log_scrollbar.set)
        
        # GUI에 로그 출력을 위한 핸들러 추가
        self.text_handler = self.TextHandler(self.log_text)
        self.text_handler.setLevel(logging.INFO)
        self.text_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        logging.getLogger().addHandler(self.text_handler)
        
        # 그리드 설정
        self.progress_bar_frame.columnconfigure(0, weight=1)
        
        # 변수 초기화
        self.ppt_path = None
        self.translation_thread = None
        self.translation_running = False
        self.start_time = 0
        self.translated_items_count = 0
        
        # 문서 분석 결과 저장
        self.total_text_elements = 0
        self.total_image_elements = 0
        self.total_elements = 0
        
        # 번역 요소 저장 (상세 프로그레스 표시용)
        self.text_elements = []  # (슬라이드_인덱스, 텍스트, 요소_타입)
        self.image_elements = [] # (슬라이드_인덱스, 이미지_패스, 이미지_크기)
        
        # 타이머 관련 변수
        self.timer_running = False
        self.timer_id = None
        self.elapsed_time = 0
        self.estimated_total_time = 0
        self.last_progress_update = 0
        
        # 초기 상태 확인
        self.check_ollama_status()
        
        # 파일 경로 변경 감지 및 버튼 색상 변경
        self.file_path_var.trace_add("write", self.on_file_path_change)
    
    # 로그 텍스트 핸들러 클래스
    class TextHandler(logging.Handler):
        def __init__(self, text_widget):
            logging.Handler.__init__(self)
            self.text_widget = text_widget
            
        def emit(self, record):
            msg = self.format(record)
            def append():
                self.text_widget.insert(tk.END, msg + '\n')
                self.text_widget.see(tk.END)
            self.text_widget.after(0, append)
    
    def on_file_path_change(self, *args):
        """파일 경로 변경 시 버튼 색상 업데이트"""
        if self.file_path_var.get().strip():
            # 파일이 선택된 경우 텍스트 색상을 검정색으로 변경
            self.start_button.config(fg="black")
        else:
            # 파일이 선택되지 않은 경우 텍스트 색상을 흰색으로 유지
            self.start_button.config(fg="white")
    
    def select_file(self):
        """파일 선택 다이얼로그 열기"""
        file_path = filedialog.askopenfilename(
            filetypes=[("PowerPoint 파일", "*.pptx"), ("모든 파일", "*.*")]
        )
        if file_path:
            self.ppt_path = file_path
            self.file_path_var.set(file_path)  # 파일 경로 표시
            self.analyze_document(file_path)
    
    def analyze_document(self, file_path):
        """문서 분석 (텍스트 및 이미지 요소 추출)"""
        self.logger.info(f"문서 분석 시작: {file_path}")
        self.status_label.config(text="문서 분석 중...")
        
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            # 분석 작업 시작
            file_name = os.path.basename(file_path)
            ppt = Presentation(file_path)
            slide_count = len(ppt.slides)
            
            # 요소 초기화
            self.text_elements = []
            self.image_elements = []
            
            total_text_count = 0
            total_image_count = 0
            total_table_cells = 0
            
            # 각 슬라이드 분석
            for slide_idx, slide in enumerate(ppt.slides):
                self.logger.debug(f"슬라이드 {slide_idx+1} 분석 중")
                
                # 텍스트 요소 분석
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # 텍스트 프레임 처리
                        if hasattr(shape, "text_frame") and shape.text.strip():
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip():
                                        self.text_elements.append({
                                            'slide_idx': slide_idx,
                                            'shape_idx': shape_idx,
                                            'type': 'text_run',
                                            'text': run.text,
                                            'translated': False
                                        })
                                        total_text_count += 1
                        
                        # 테이블 처리
                        if hasattr(shape, "table"):
                            table = shape.table
                            for row_idx, row in enumerate(table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    if cell.text.strip():
                                        self.text_elements.append({
                                            'slide_idx': slide_idx,
                                            'shape_idx': shape_idx,
                                            'type': 'table_cell',
                                            'row_idx': row_idx,
                                            'col_idx': col_idx,
                                            'text': cell.text,
                                            'translated': False
                                        })
                                        total_text_count += 1
                                        total_table_cells += 1
                        
                        # 이미지 처리
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image = shape.image
                            image_bytes = image.blob
                            image_size = len(image_bytes)
                            
                            self.image_elements.append({
                                'slide_idx': slide_idx,
                                'shape_idx': shape_idx,
                                'type': 'image',
                                'size': image_size,
                                'translated': False
                            })
                            total_image_count += 1
                            
                    except Exception as e:
                        self.logger.error(f"요소 분석 오류 (슬라이드 {slide_idx+1}, 요소 {shape_idx}): {str(e)}")
            
            # 총 요소 수 계산
            self.total_text_elements = total_text_count
            self.total_image_elements = total_image_count
            self.total_elements = total_text_count + total_image_count
            
            # 정보 표시
            self.file_name_label.config(text=f"파일 이름: {file_name}")
            self.slide_count_label.config(text=f"슬라이드 수: {slide_count}")
            self.text_count_label.config(text=f"텍스트 요소 수: {total_text_count} (테이블 셀: {total_table_cells})")
            self.image_count_label.config(text=f"이미지 요소 수: {total_image_count}")
            self.total_elements_label.config(text=f"총 번역 요소: {self.total_elements}")
            
            # 분석 결과 로깅
            self.logger.info(f"문서 분석 완료: 슬라이드 {slide_count}개, 텍스트 요소 {total_text_count}개, 이미지 {total_image_count}개")
            
            # 번역 상태 초기화
            self.translated_items_count = 0
            self.translated_items_label.config(text=f"번역된 요소: 0/{self.total_elements}")
            self.remaining_items_label.config(text=f"남은 요소: {self.total_elements}")
            
            self.status_label.config(text="문서 분석 완료")
            
        except Exception as e:
            self.logger.exception(f"문서 분석 오류: {str(e)}")
            self.file_name_label.config(text=f"파일 이름: {os.path.basename(file_path)}")
            self.slide_count_label.config(text="슬라이드 수: 오류 발생")
            self.text_count_label.config(text="텍스트 요소 수: 오류 발생")
            self.image_count_label.config(text="이미지 요소 수: 오류 발생")
            self.total_elements_label.config(text="총 번역 요소: 오류 발생")
            self.status_label.config(text=f"문서 분석 오류: {str(e)}")
            messagebox.showerror("오류", f"문서 분석 중 오류가 발생했습니다: {str(e)}")
    
    def check_ollama_status(self):
        """Ollama 상태 확인"""
        # 설치 확인
        installed = self.is_ollama_installed()
        self.ollama_installed_label.config(
            text=f"설치 상태: {'설치됨' if installed else '설치되지 않음'}",
            fg="green" if installed else "red"
        )
        
        if not installed:
            self.show_ollama_install_guide()
            return False
        
        # 실행 상태 확인
        running, port = self.is_ollama_running()
        self.ollama_running_label.config(
            text=f"실행 상태: {'실행 중' if running else '실행되지 않음'}",
            fg="green" if running else "red"
        )
        
        # Ollama가 설치되어 있지만 실행 중이 아닌 경우 자동 실행
        if installed and not running:
            self.start_ollama()
            running, port = self.is_ollama_running()
        
        # 포트 정보 표시
        self.ollama_port_label.config(
            text=f"Ollama 포트: {port if running else '없음'}"
        )
        
        # URL 업데이트
        if running and port:
            self.url_var.set(f"http://localhost:{port}")
            
            # 최초 실행 시에만 모델 목록 업데이트
            if not self.models_initialized:
                text_models, vision_models = self.update_models_list()
                
                # 필요한 모델이 설치되어 있는지 확인
                if not text_models and not vision_models:
                    self.prompt_install_base_models()
                elif not vision_models:
                    self.prompt_install_vision_model()
                
                self.models_initialized = True
        
        return installed and running
    
    def show_ollama_install_guide(self):
        """Ollama 설치 가이드 표시"""
        response = messagebox.askquestion(
            "Ollama 설치 필요",
            "Ollama가 설치되어 있지 않습니다. Ollama 설치 페이지로 이동하시겠습니까?",
            icon='warning'
        )
        
        if response == 'yes':
            webbrowser.open("https://ollama.com/download")
    
    def start_ollama(self):
        """Ollama 시작"""
        try:
            self.status_label.config(text="Ollama 시작 중...")
            self.logger.info("Ollama 시작 시도")
            
            if platform.system() == "Windows":
                subprocess.Popen(["ollama", "serve"], shell=True, creationflags=subprocess.DETACHED_PROCESS)
            else:  # macOS 또는 Linux
                subprocess.Popen(["ollama", "serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            
            # Ollama가 시작될 때까지 대기
            for _ in range(10):  # 최대 10초 대기
                time.sleep(1)
                running, _ = self.is_ollama_running()
                if running:
                    self.status_label.config(text="Ollama 시작됨")
                    self.logger.info("Ollama 시작 성공")
                    return True
            
            self.status_label.config(text="Ollama 시작 실패")
            self.logger.warning("Ollama 시작 실패: 시간 초과")
            return False
        except Exception as e:
            self.status_label.config(text=f"Ollama 시작 오류: {str(e)}")
            self.logger.error(f"Ollama 시작 오류: {str(e)}")
            return False
    
    def prompt_install_base_models(self):
        """기본 모델 설치 권유"""
        response = messagebox.askquestion(
            "기본 모델 설치",
            "사용 가능한 모델이 없습니다. 기본 모델(gemma3:12b, llava)을 설치하시겠습니까?",
            icon='info'
        )
        
        if response == 'yes':
            self.status_label.config(text="기본 모델 설치 중...")
            threading.Thread(target=self.install_base_models, daemon=True).start()
    
    def prompt_install_vision_model(self):
        """Vision 모델 설치 권유"""
        response = messagebox.askquestion(
            "Vision 모델 설치",
            "이미지 처리를 위한 Vision 모델이 없습니다. llava 모델을 설치하시겠습니까?",
            icon='info'
        )
        
        if response == 'yes':
            self.status_label.config(text="llava 모델 설치 중...")
            threading.Thread(target=lambda: self.install_model("llava"), daemon=True).start()
    
    def install_base_models(self):
        """기본 모델 설치"""
        try:
            # Gemma 3:12b 설치
            self.install_model("gemma3:12b")
            
            # LLaVA 설치
            self.install_model("llava")
            
            # 모델 목록 업데이트
            self.update_models_list()
            
            self.status_label.config(text="모델 설치 완료")
            messagebox.showinfo("설치 완료", "기본 모델 설치가 완료되었습니다.")
        except Exception as e:
            self.status_label.config(text=f"모델 설치 오류: {str(e)}")
            messagebox.showerror("설치 오류", f"모델 설치 중 오류가 발생했습니다: {str(e)}")
    
    def install_model(self, model_name):
        """모델 설치"""
        try:
            self.status_label.config(text=f"{model_name} 설치 중...")
            self.logger.info(f"{model_name} 설치 시작")
            
            if platform.system() == "Windows":
                subprocess.run(["ollama", "pull", model_name], shell=True, check=True)
            else:  # macOS 또는 Linux
                subprocess.run(["ollama", "pull", model_name], check=True)
            
            self.logger.info(f"{model_name} 설치 완료")
            return True
        except Exception as e:
            self.logger.error(f"{model_name} 설치 오류: {e}")
            return False
    
    def update_models_list(self):
        """설치된 모델 목록 가져오기 및 UI 업데이트"""
        try:
            # 현재 선택된 모델 저장
            current_text_model = self.text_model_var.get()
            current_vision_model = self.vision_model_var.get()
            
            text_models = []
            vision_models = []
            
            # API로 모델 목록 가져오기
            self.logger.debug(f"모델 목록 가져오기: {self.url_var.get()}/api/tags")
            response = requests.get(f"{self.url_var.get()}/api/tags", timeout=5)
            if response.status_code == 200:
                models_data = response.json()
                if 'models' in models_data:
                    all_models = [model['name'] for model in models_data['models']]
                    self.logger.debug(f"발견된 모델: {all_models}")
                    
                    # 모델 분류
                    for model in all_models:
                        if any(vm in model.lower() for vm in ["llava", "bakllava", "vision", "vl", "multimodal"]):
                            vision_models.append(model)
                        else:
                            text_models.append(model)
                    
                    # UI 업데이트
                    if vision_models:
                        self.vision_model_combo['values'] = vision_models
                        if current_vision_model in vision_models:
                            self.vision_model_var.set(current_vision_model)
                        else:
                            self.vision_model_var.set(vision_models[0])
                    else:
                        self.vision_model_combo['values'] = ["Vision 모델 없음"]
                        self.vision_model_var.set("Vision 모델 없음")
                    
                    if text_models:
                        self.text_model_combo['values'] = text_models
                        if current_text_model in text_models:
                            self.text_model_var.set(current_text_model)
                        elif "gemma3:12b" in text_models:
                            self.text_model_var.set("gemma3:12b")
                        else:
                            self.text_model_var.set(text_models[0])
                    else:
                        self.text_model_combo['values'] = ["Text 모델 없음"]
                        self.text_model_var.set("Text 모델 없음")
                    
                    return text_models, vision_models
            
            # API 방식이 실패한 경우 명령행 방식 시도
            if platform.system() != "Windows":
                self.logger.debug("명령행으로 모델 목록 가져오기 시도")
                result = subprocess.run(["ollama", "list"], capture_output=True, text=True)
                if result.returncode == 0:
                    lines = result.stdout.strip().split('\n')
                    if len(lines) > 1:  # 헤더 제외
                        all_models = []
                        for line in lines[1:]:
                            parts = line.split()
                            if parts:
                                all_models.append(parts[0])
                        
                        self.logger.debug(f"명령행으로 발견된 모델: {all_models}")
                        
                        # 모델 분류
                        for model in all_models:
                            if any(vm in model.lower() for vm in ["llava", "bakllava", "vision", "vl", "multimodal"]):
                                vision_models.append(model)
                            else:
                                text_models.append(model)
                        
                        # UI 업데이트
                        if vision_models:
                            self.vision_model_combo['values'] = vision_models
                            if current_vision_model in vision_models:
                                self.vision_model_var.set(current_vision_model)
                            else:
                                self.vision_model_var.set(vision_models[0])
                        else:
                            self.vision_model_combo['values'] = ["Vision 모델 없음"]
                            self.vision_model_var.set("Vision 모델 없음")
                        
                        if text_models:
                            self.text_model_combo['values'] = text_models
                            if current_text_model in text_models:
                                self.text_model_var.set(current_text_model)
                            elif "gemma3:12b" in text_models:
                                self.text_model_var.set("gemma3:12b")
                            else:
                                self.text_model_var.set(text_models[0])
                        else:
                            self.text_model_combo['values'] = ["Text 모델 없음"]
                            self.text_model_var.set("Text 모델 없음")
                        
                        return text_models, vision_models
            
            # 모두 실패한 경우 빈 목록으로 설정
            self.logger.warning("모델 목록을 가져올 수 없음")
            self.vision_model_combo['values'] = ["Vision 모델 없음"]
            self.vision_model_var.set("Vision 모델 없음")
            self.text_model_combo['values'] = ["Text 모델 없음"]
            self.text_model_var.set("Text 모델 없음")
            
            return [], []
            
        except Exception as e:
            self.logger.exception(f"모델 목록 가져오기 오류: {e}")
            self.vision_model_combo['values'] = ["Vision 모델 없음"]
            self.vision_model_var.set("Vision 모델 없음")
            self.text_model_combo['values'] = ["Text 모델 없음"]
            self.text_model_var.set("Text 모델 없음")
            
            return [], []
    
    def is_ollama_installed(self):
        """Ollama 설치 여부 확인"""
        try:
            system = platform.system()
            if system == "Windows":
                return os.path.exists("C:\\Program Files\\Ollama\\ollama.exe") or \
                       os.path.exists(os.path.expanduser("~\\AppData\\Local\\Ollama\\ollama.exe"))
            elif system == "Darwin":  # macOS
                return os.path.exists("/usr/local/bin/ollama") or \
                       os.path.exists("/opt/homebrew/bin/ollama")
            elif system == "Linux":
                result = subprocess.run(["which", "ollama"], capture_output=True, text=True)
                return result.returncode == 0
            return False
        except Exception as e:
            self.logger.error(f"Ollama 설치 확인 오류: {e}")
            return False
    
    def is_ollama_running(self):
        """Ollama 실행 상태 및 포트 확인"""
        try:
            # API 호출 시도
            try:
                self.logger.debug(f"Ollama API 호출: {self.url_var.get()}/api/tags")
                response = requests.get(self.url_var.get() + "/api/tags", timeout=2)
                if response.status_code == 200:
                    port = self.url_var.get().split(':')[-1]
                    self.logger.debug(f"Ollama 실행 중: 포트 {port}")
                    return True, port
            except Exception as e:
                self.logger.debug(f"API 호출 실패: {e}")
                pass
            
            # 프로세스 확인
            for proc in psutil.process_iter(['pid', 'name']):
                if 'ollama' in proc.info['name'].lower():
                    self.logger.debug("Ollama 프로세스 발견")
                    return True, 11434  # 기본 포트
            
            self.logger.debug("Ollama가 실행 중이 아님")
            return False, None
        except Exception as e:
            self.logger.exception(f"Ollama 실행 상태 확인 오류: {e}")
            return False, None
    
    def extract_text_with_vision_model(self, image_path, model_name):
        """Vision 모델을 사용하여 이미지에서 텍스트 추출"""
        self.logger.info(f"이미지 텍스트 추출 시작: {image_path}")
        
        if "Vision 모델 없음" in model_name:
            self.logger.warning("Vision 모델이 없어 텍스트 추출할 수 없음")
            return ""
        
        try:
            # 이미지 크기 확인 및 리사이징
            img = Image.open(image_path)
            img_size = os.path.getsize(image_path)
            self.logger.info(f"원본 이미지 크기: {img.width}x{img.height}, {img_size} 바이트")
            
            # 큰 이미지 리사이징 (10MB 초과)
            if img_size > 10 * 1024 * 1024:  # 10MB
                self.logger.info("큰 이미지 리사이징")
                max_size = 1024  # 최대 1024px 크기로 제한
                if img.width > max_size or img.height > max_size:
                    # 비율 유지하면서 리사이징
                    ratio = min(max_size / img.width, max_size / img.height)
                    new_width = int(img.width * ratio)
                    new_height = int(img.height * ratio)
                    img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                    
                    # 리사이징된 이미지 저장
                    resized_path = f"{os.path.splitext(image_path)[0]}_resized.png"
                    img.save(resized_path)
                    image_path = resized_path
                    self.logger.info(f"이미지 리사이징 완료: {new_width}x{new_height}")
            
            # 이미지를 base64로 인코딩
            with open(image_path, "rb") as img_file:
                img_data = img_file.read()
                self.logger.info(f"처리할 이미지 크기: {len(img_data)} 바이트")
                img_base64 = base64.b64encode(img_data).decode('utf-8')
            
            # Ollama API 호출
            self.logger.info(f"Vision API 호출 시작: {model_name}")
            start_time = time.time()
            
            response = requests.post(
                f"{self.url_var.get()}/api/generate",
                json={
                    "model": model_name,
                    "prompt": "Extract all visible text from this image. Only return the text, nothing else.",
                    "images": [img_base64]
                },
                timeout=60  # 60초 타임아웃 설정
            )
            
            elapsed = time.time() - start_time
            self.logger.info(f"Vision API 응답 수신: {elapsed:.2f}초 소요")
            
            if response.status_code == 200:
                result = response.json()
                if 'response' in result:
                    extracted_text = result['response'].strip()
                    self.logger.info(f"추출된 텍스트 길이: {len(extracted_text)} 글자, 내용: '{extracted_text[:50]}...'")
                    return extracted_text
                else:
                    self.logger.error(f"API 응답에 'response' 필드 없음: {result}")
                    return ""
            else:
                self.logger.error(f"Vision API 오류 (HTTP {response.status_code}): {response.text}")
                return ""
        except requests.exceptions.Timeout:
            self.logger.error("Vision API 호출 타임아웃")
            return ""
        except json.JSONDecodeError as e:
            self.logger.error(f"JSON 파싱 오류: {e} - 응답: {response.text[:100]}")
            return ""
        except Exception as e:
            self.logger.exception(f"텍스트 추출 오류: {str(e)}")
            return ""
    
    def translate_text_with_ollama(self, text, source_lang, target_lang, model, url):
        """Text 모델을 사용하여 텍스트 번역"""
        if not text or text.isspace() or "Text 모델 없음" in model:
            return text
        
        self.logger.debug(f"번역 시작: '{text[:50]}...'")
        
        # 번역 프롬프트
        prompt = f"You are a translator. Your role is to accurately translate the given {source_lang} text into {target_lang}. Do not provide any explanations, only the translated result. : {text}"
        
        try:
            start_time = time.time()
            response = requests.post(
                f"{url}/api/generate",
                json={
                    "model": model,
                    "prompt": prompt,
                    "stream": False
                },
                timeout=30  # 30초 타임아웃
            )
            
            elapsed = time.time() - start_time
            self.logger.debug(f"번역 API 응답 시간: {elapsed:.2f}초")
            
            if response.status_code == 200:
                try:
                    result = response.json()
                    translated_text = result.get("response", "").strip()
                    
                    self.logger.info(f"번역 완료: '{text[:30]}...' → '{translated_text[:30]}...'")
                    
                    self.translated_items_count += 1
                    self.translated_items_label.config(text=f"번역된 요소: {self.translated_items_count}/{self.total_elements}")
                    self.remaining_items_label.config(text=f"남은 요소: {self.total_elements - self.translated_items_count}")
                    self.root.update_idletasks()
                    return translated_text
                except json.JSONDecodeError as e:
                    self.logger.error(f"JSON 파싱 오류: {e} - 응답: {response.text[:100]}")
                    return text
            else:
                self.logger.error(f"번역 API 오류 (HTTP {response.status_code}): {response.text[:100]}")
                return text
        except requests.exceptions.Timeout:
            self.logger.error("번역 API 타임아웃")
            return text
        except Exception as e:
            self.logger.exception(f"번역 오류: {str(e)}")
            return text
    
    def translate_image(self, image_path, source_lang, target_lang, vision_model, text_model, url):
        """이미지 내 텍스트 추출 및 번역"""
        self.logger.info(f"이미지 번역 시작: {image_path}")
        
        try:
            # 원본 이미지 로드
            img = cv2.imread(image_path)
            if img is None:
                self.logger.error(f"이미지 로드 실패: {image_path}")
                return image_path
            
            # Vision 모델로 텍스트 추출
            extracted_text = self.extract_text_with_vision_model(image_path, vision_model)
            
            if not extracted_text:
                self.logger.warning("이미지에서 텍스트를 추출할 수 없음")
                return image_path
            
            # 추출된 텍스트 번역
            translated_text = self.translate_text_with_ollama(
                extracted_text, source_lang, target_lang, text_model, url
            )
            
            if not translated_text or translated_text == extracted_text:
                self.logger.warning("텍스트 번역 결과가 없거나 원본과 동일")
                return image_path
            
            # 결과를 이미지에 삽입 (오버레이 텍스트 방식)
            overlay = img.copy()
            h, w = img.shape[:2]
            
            # 배경색 설정 (흰색 반투명)
            bg_color = (255, 255, 255)
            alpha = 0.7  # 투명도
            
            # 텍스트 분할 (여러 줄로)
            lines = []
            words = translated_text.split()
            line = ""
            max_line_length = 50  # 한 줄에 최대 50자
            
            for word in words:
                test_line = line + " " + word if line else word
                if len(test_line) <= max_line_length:
                    line = test_line
                else:
                    lines.append(line)
                    line = word
            
            if line:
                lines.append(line)
            
            # 텍스트 크기 계산
            font = cv2.FONT_HERSHEY_SIMPLEX
            font_scale = 0.7
            thickness = 2
            line_height = 30
            
            text_height = len(lines) * line_height
            text_width = 0
            
            for line in lines:
                (line_width, _), _ = cv2.getTextSize(line, font, font_scale, thickness)
                text_width = max(text_width, line_width)
            
            # 텍스트 위치 설정 (이미지 하단에 배치)
            text_x = 10
            text_y = h - text_height - 10
            
            # 텍스트 배경 생성
            cv2.rectangle(overlay, (text_x - 5, text_y - 15), 
                         (text_x + text_width + 5, text_y + text_height + 5), 
                         bg_color, -1)
            
            # 투명도 적용
            cv2.addWeighted(overlay, alpha, img, 1 - alpha, 0, img)
            
            # 텍스트 추가
            for i, line in enumerate(lines):
                y = text_y + i * line_height
                cv2.putText(img, line, (text_x, y), font, font_scale, (0, 0, 0), thickness)
            
            # 결과 이미지 저장
            output_path = f"translated_{os.path.basename(image_path)}"
            cv2.imwrite(output_path, img)
            self.logger.info(f"번역된 이미지 저장: {output_path}")
            
            return output_path
            
        except Exception as e:
            self.logger.exception(f"이미지 번역 오류: {str(e)}")
            return image_path
    
    def format_time(self, seconds):
        """초를 분:초 형식으로 변환"""
        minutes = int(seconds // 60)
        seconds = int(seconds % 60)
        return f"{minutes}:{seconds:02d}"
    
    def update_timer(self):
        """타이머 업데이트 함수"""
        if self.timer_running:
            # 경과 시간 계산
            self.elapsed_time = time.time() - self.start_time
            
            # 진행율 가져오기
            progress = self.progress_var.get()
            
            # 경과 시간만 업데이트 (예상 시간은 프로그레스 업데이트 시에만 변경)
            time_text = self.progress_label.cget("text")
            if "계산 중" in time_text or progress < 1:
                self.progress_label.config(
                    text=f"{progress:.1f}% ({self.format_time(self.elapsed_time)} / 계산 중)"
                )
            
            # 1초 후 다시 호출
            self.timer_id = self.root.after(1000, self.update_timer)
    
    def update_progress(self, current, total):
        """진행 상황 업데이트"""
        if total == 0:
            return
        
        # 진행율 계산
        progress = (current / total) * 100
        self.progress_var.set(progress)
        
        # 예상 시간은 프로그레스가 갱신될 때만 업데이트
        # 경과 시간 계산
        self.elapsed_time = time.time() - self.start_time
        
        # 예상 남은 시간 계산
        if progress > 5:  # 최소 5% 이상 진행되었을 때만 예상 시간 계산
            self.estimated_total_time = self.elapsed_time * (100 / progress)
            estimated_remaining_time = self.estimated_total_time - self.elapsed_time
            
            # 진행 상황 표시 업데이트
            self.progress_label.config(
                text=f"{progress:.1f}% ({self.format_time(self.elapsed_time)} / {self.format_time(self.elapsed_time + estimated_remaining_time)})"
            )
        else:
            # 진행률이 낮을 때는 "계산 중" 표시
            self.progress_label.config(
                text=f"{progress:.1f}% ({self.format_time(self.elapsed_time)} / 계산 중)"
            )
        
        self.root.update_idletasks()
    
    def start_translation(self):
        """번역 프로세스 시작"""
        if not self.file_path_var.get():
            messagebox.showerror("오류", "파일을 먼저 선택해주세요.")
            return
        
        self.ppt_path = self.file_path_var.get()
        
        # 파일 존재 확인
        if not os.path.exists(self.ppt_path):
            messagebox.showerror("오류", "선택한 파일이 존재하지 않습니다.")
            return
        
        # Ollama 상태 확인
        installed = self.is_ollama_installed()
        running, port = self.is_ollama_running()
        if not (installed and running):
            messagebox.showerror("오류", "Ollama가 설치되어 있지 않거나 실행 중이 아닙니다.")
            return
        
        # 모델 확인
        vision_model = self.vision_model_var.get()
        text_model = self.text_model_var.get()
        
        if "Vision 모델 없음" in vision_model:
            response = messagebox.askquestion(
                "Vision 모델 없음",
                "이미지 처리를 위한 Vision 모델이 없습니다. Vision 모델을 설치한 후 다시 시도하시겠습니까?",
                icon='warning'
            )
            if response == 'yes':
                self.prompt_install_vision_model()
            return
        
        if "Text 모델 없음" in text_model:
            response = messagebox.askquestion(
                "Text 모델 없음",
                "텍스트 번역을 위한 모델이 없습니다. Text 모델을.설치한 후 다시 시도하시겠습니까?",
                icon='warning'
            )
            if response == 'yes':
                self.status_label.config(text="gemma3:12b 모델 설치 중...")
                threading.Thread(target=lambda: self.install_model("gemma3:12b"), daemon=True).start()
            return
        
        # 이미 번역 중인지 확인
        if self.translation_running:
            messagebox.showinfo("알림", "이미 번역이 진행 중입니다.")
            return
        
        # 파일 분석이 필요한 경우
        if self.total_elements == 0:
            self.analyze_document(self.ppt_path)
            if self.total_elements == 0:
                messagebox.showerror("오류", "번역할 요소가 없습니다.")
                return
        
        # 진행 상태 초기화
        self.progress_var.set(0)
        self.progress_label.config(text="0% (0:00 / 계산 중)")
        self.current_slide_label.config(text="현재 슬라이드: -")
        self.current_task_label.config(text="현재 작업: -")
        self.translated_items_count = 0
        self.translated_items_label.config(text=f"번역된 요소: 0/{self.total_elements}")
        self.remaining_items_label.config(text=f"남은 요소: {self.total_elements}")
        
        # 타이머 초기화
        self.elapsed_time = 0
        self.estimated_total_time = 0
        self.last_progress_update = 0
        
        # 버튼 상태 변경
        self.start_button.config(state=tk.DISABLED)
        self.stop_button.config(state=tk.NORMAL)
        
        # 번역 스레드 시작
        self.translation_running = True
        self.start_time = time.time()
        self.translation_thread = threading.Thread(target=self.translation_process)
        self.translation_thread.daemon = True
        self.translation_thread.start()
        
        # 타이머 시작
        self.timer_running = True
        self.update_timer()
    
    def stop_translation(self):
        """번역 프로세스 중지"""
        self.translation_running = False
        self.timer_running = False  # 타이머 중지
        self.status_label.config(text="번역 중지 중...")
        self.logger.info("사용자에 의한 번역 중지")
    
    def translation_process(self):
        """번역 프로세스 실행"""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            self.logger.info(f"번역 프로세스 시작: {self.ppt_path}")
            
            # 옵션 가져오기
            source_lang = self.source_lang.get()
            target_lang = self.target_lang.get()
            vision_model = self.vision_model_var.get()
            text_model = self.text_model_var.get()
            url = self.url_var.get()
            
            self.logger.info(f"번역 설정: {source_lang} → {target_lang}, Vision: {vision_model}, Text: {text_model}")
            
            # 파워포인트 파일 열기
            ppt = Presentation(self.ppt_path)
            
            # 임시 폴더 생성
            with tempfile.TemporaryDirectory() as temp_dir:
                self.logger.info(f"임시 폴더 생성: {temp_dir}")
                
                # 프로그레스 초기화
                processed_items = 0
                
                # 1. 텍스트 요소 번역
                self.current_task_label.config(text="현재 작업: 텍스트 번역")
                self.status_label.config(text="텍스트 요소 번역 중...")
                self.logger.info("텍스트 요소 번역 시작")
                
                # 텍스트 요소 전체 개수 가져오기
                total_text_elements = len(self.text_elements)
                
                for idx, text_element in enumerate(self.text_elements):
                    if not self.translation_running:
                        raise Exception("사용자에 의해 번역이 중지되었습니다.")
                    
                    slide_idx = text_element['slide_idx']
                    slide = ppt.slides[slide_idx]
                    
                    self.current_slide_label.config(text=f"현재 슬라이드: {slide_idx+1}/{len(ppt.slides)}")
                    
                    try:
                        if text_element['type'] == 'text_run':
                            self.logger.info(f"텍스트 번역 (슬라이드 {slide_idx+1}, 요소 {text_element['shape_idx']}): '{text_element['text'][:30]}...'")
                            
                            # 해당 텍스트 요소 찾기
                            shape = slide.shapes[text_element['shape_idx']]
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip() == text_element['text']:
                                        # 번역
                                        translated_text = self.translate_text_with_ollama(
                                            run.text, source_lang, target_lang, text_model, url
                                        )
                                        run.text = translated_text
                                        text_element['translated'] = True
                                        break
                        
                        elif text_element['type'] == 'table_cell':
                            self.logger.info(f"테이블 셀 번역 (슬라이드 {slide_idx+1}, 테이블 {text_element['shape_idx']}, 행 {text_element['row_idx']}, 열 {text_element['col_idx']}): '{text_element['text'][:30]}...'")
                            
                            # 해당 테이블 찾기
                            shape = slide.shapes[text_element['shape_idx']]
                            if hasattr(shape, "table"):
                                table = shape.table
                                cell = table.rows[text_element['row_idx']].cells[text_element['col_idx']]
                                
                                # 번역
                                if cell.text.strip() == text_element['text']:
                                    translated_text = self.translate_text_with_ollama(
                                        cell.text, source_lang, target_lang, text_model, url
                                    )
                                    
                                    # 텍스트 설정 (레퍼런스 이슈 방지)
                                    text_frame = cell.text_frame
                                    # 모든 기존 단락 제거
                                    while len(text_frame.paragraphs) > 1:
                                        tr_element = text_frame._txBody.remove(text_frame._txBody[1])
                                    
                                    # 첫 번째 단락에 새 텍스트 설정
                                    text_frame.paragraphs[0].text = translated_text
                                    
                                    text_element['translated'] = True
                    
                    except Exception as e:
                        self.logger.error(f"텍스트 번역 오류 (요소 {idx+1}/{total_text_elements}): {str(e)}")
                        self.logger.error(traceback.format_exc())
                    
                    # 진행 상황 업데이트
                    processed_items += 1
                    self.update_progress(processed_items, self.total_elements)
                
                # 2. 이미지 요소 번역
                self.current_task_label.config(text="현재 작업: 이미지 번역")
                self.status_label.config(text="이미지 요소 번역 중...")
                self.logger.info("이미지 요소 번역 시작")
                
                # 이미지 요소 전체 개수 가져오기
                total_image_elements = len(self.image_elements)
                
                for idx, image_element in enumerate(self.image_elements):
                    if not self.translation_running:
                        raise Exception("사용자에 의해 번역이 중지되었습니다.")
                    
                    slide_idx = image_element['slide_idx']
                    slide = ppt.slides[slide_idx]
                    
                    self.current_slide_label.config(text=f"현재 슬라이드: {slide_idx+1}/{len(ppt.slides)}")
                    
                    try:
                        self.logger.info(f"이미지 번역 (슬라이드 {slide_idx+1}, 요소 {image_element['shape_idx']})")
                        
                        # 해당 이미지 찾기
                        shape = slide.shapes[image_element['shape_idx']]
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # 이미지 추출 및 저장
                            image = shape.image
                            image_bytes = image.blob
                            
                            temp_image_path = os.path.join(temp_dir, f"slide_{slide_idx}_image_{image_element['shape_idx']}.png")
                            with open(temp_image_path, "wb") as f:
                                f.write(image_bytes)
                            
                            self.logger.info(f"이미지 저장: {temp_image_path} ({len(image_bytes)} 바이트)")
                            
                            # 이미지 번역
                            translated_image_path = self.translate_image(
                                temp_image_path, source_lang, target_lang, vision_model, text_model, url
                            )
                            
                            # 번역된 이미지로 교체
                            try:
                                if os.path.exists(translated_image_path) and translated_image_path != temp_image_path:
                                    # 이미지 교체 방법 1: 원본 이미지와 크기, 형식이 같은 경우
                                    with open(translated_image_path, "rb") as f:
                                        new_image_bytes = f.read()
                                        
                                    # 이미지 교체 방법 (Python-pptx 제한 우회)
                                    # 1. 새 이미지 슬라이드에 추가
                                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                    pic = slide.shapes.add_picture(translated_image_path, left, top, width, height)
                                    
                                    # 2. 기존 이미지와 동일한 위치/크기 설정
                                    pic.left, pic.top, pic.width, pic.height = left, top, width, height
                                    
                                    # 3. 기존 이미지 제거 (실제 PowerPoint에서는 불가능할 수 있음)
                                    # 대신 투명하게 만들거나 배경색으로 덮을 수 있음
                                    image_element['translated'] = True
                                    
                                    self.logger.info("이미지 교체 완료")
                            except Exception as e:
                                self.logger.error(f"이미지 교체 오류: {str(e)}")
                                self.logger.error(traceback.format_exc())
                    
                    except Exception as e:
                        self.logger.error(f"이미지 번역 오류 (요소 {idx+1}/{total_image_elements}): {str(e)}")
                        self.logger.error(traceback.format_exc())
                    
                    # 진행 상황 업데이트
                    processed_items += 1
                    self.update_progress(processed_items, self.total_elements)
                
                # 번역된 파일 저장
                output_path = os.path.splitext(self.ppt_path)[0] + "_translated.pptx"
                self.logger.info(f"번역된 파일 저장: {output_path}")
                ppt.save(output_path)
                
                # 타이머 중지
                self.timer_running = False
                
                # 완료 메시지
                elapsed_time = time.time() - self.start_time
                self.progress_label.config(text=f"100% (총 소요시간: {self.format_time(elapsed_time)})")
                self.status_label.config(text=f"번역 완료! 파일 저장됨: {output_path}")
                self.logger.info(f"번역 완료: 소요시간 {elapsed_time:.2f}초")
                messagebox.showinfo("완료", f"번역이 완료되었습니다.\n파일 저장 위치: {output_path}")
            
        except Exception as e:
            # 오류 발생 시 타이머 중지
            self.timer_running = False
            self.logger.exception(f"번역 프로세스 오류: {str(e)}")
            print(f"번역 오류: {e}")
            self.status_label.config(text=f"번역 오류: {str(e)}")
            messagebox.showerror("오류", f"번역 중 오류가 발생했습니다: {str(e)}")
        
        finally:
            self.translation_running = False
            self.timer_running = False
            self.start_button.config(state=tk.NORMAL)
            self.stop_button.config(state=tk.DISABLED)

if __name__ == "__main__":
    root = tk.Tk()
    app = PowerPointTranslator(root)
    root.mainloop()