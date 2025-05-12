# services/document_analyzer.py
import os
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class DocumentAnalyzer:
    def analyze_ppt(self, file_path):
        """PPT 파일 분석 (텍스트 요소를 문단 단위로 추출)"""
        logger.info(f"문서 분석 시작: {file_path}")
        
        try:
            # 분석 작업 시작
            file_name = os.path.basename(file_path)
            ppt = Presentation(file_path)
            slide_count = len(ppt.slides)
            
            # 요소 초기화
            text_elements = []
            image_elements = []
            
            total_text_count = 0
            total_image_count = 0
            total_table_cells = 0
            
            # 각 슬라이드 분석
            for slide_idx, slide in enumerate(ppt.slides):
                logger.debug(f"슬라이드 {slide_idx+1} 분석 중")
                self._analyze_slide(slide, slide_idx, text_elements, image_elements, 
                                   total_text_count, total_image_count, total_table_cells)
            
            # 총 요소 수 계산
            total_elements = len(text_elements) + len(image_elements)
            total_text_count = len(text_elements)
            total_image_count = len(image_elements)
            
            # 결과 저장
            result = {
                'file_name': file_name,
                'slide_count': slide_count,
                'text_elements': text_elements,
                'image_elements': image_elements,
                'total_text_count': total_text_count,
                'total_image_count': total_image_count,
                'total_table_cells': total_table_cells,
                'total_elements': total_elements
            }
            
            logger.info(f"문서 분석 완료: 슬라이드 {slide_count}개, 텍스트 요소 {total_text_count}개, 이미지 {total_image_count}개")
            return result
            
        except Exception as e:
            logger.exception(f"문서 분석 오류: {str(e)}")
            raise
            
    def _analyze_slide(self, slide, slide_idx, text_elements, image_elements, 
                      total_text_count, total_image_count, total_table_cells):
        """개별 슬라이드 분석"""
        # 각 요소 분석
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # 텍스트 프레임 처리
                if hasattr(shape, "text_frame") and shape.text.strip():
                    self._process_text_frame(shape, slide_idx, shape_idx, text_elements)
                
                # 테이블 처리
                if hasattr(shape, "table"):
                    self._process_table(shape, slide_idx, shape_idx, text_elements, total_table_cells)
                
                # 이미지 처리
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._process_image(shape, slide_idx, shape_idx, image_elements)
                    
            except Exception as e:
                logger.error(f"요소 분석 오류 (슬라이드 {slide_idx+1}, 요소 {shape_idx}): {str(e)}")
    
    def _process_text_frame(self, shape, slide_idx, shape_idx, text_elements):
        """텍스트 프레임 처리"""
        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            if paragraph.text.strip():
                text_elements.append({
                    'slide_idx': slide_idx,
                    'shape_idx': shape_idx,
                    'para_idx': para_idx,
                    'type': 'paragraph',
                    'text': paragraph.text.strip(),
                    'translated': False
                })
    
    def _process_table(self, shape, slide_idx, shape_idx, text_elements, total_table_cells):
        """테이블 처리"""
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text.strip():
                    text_elements.append({
                        'slide_idx': slide_idx,
                        'shape_idx': shape_idx,
                        'type': 'table_cell',
                        'row_idx': row_idx,
                        'col_idx': col_idx,
                        'text': cell.text.strip(),
                        'translated': False
                    })
                    total_table_cells += 1
    
    def _process_image(self, shape, slide_idx, shape_idx, image_elements):
        """이미지 처리"""
        image = shape.image
        image_bytes = image.blob
        image_size = len(image_bytes)
        
        image_elements.append({
            'slide_idx': slide_idx,
            'shape_idx': shape_idx,
            'type': 'image',
            'size': image_size,
            'translated': False
        })