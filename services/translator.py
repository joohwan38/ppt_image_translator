import os
import tempfile
import time
import logging
import traceback
import base64
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from utils.image_utils import resize_image_if_needed, encode_image_to_base64, overlay_text_on_image

logger = logging.getLogger(__name__)

class TranslationService:
    def __init__(self, ollama_service):
        self.ollama_service = ollama_service

    def is_numeric_text(self, text):
        """숫자와 관련된 텍스트 감지 (숫자, 단위, 특수 형식 등)"""
        # 공백 제거
        text = text.strip()
        
        # 빈 문자열이면 False 반환
        if not text:
            return False
        
        # 순수 숫자 확인 (쉼표 포함)
        numeric_only = text.replace(',', '').replace('.', '')
        if numeric_only.isdigit():
            return True
        
        # 퍼센트 표시 감지 (재귀 호출 제거)
        if text.endswith('%'):
            # %를 제외한 부분이 숫자인지 확인
            text_without_percent = text[:-1].strip()
            try:
                float(text_without_percent.replace(',', ''))
                return True
            except ValueError:
                pass
        
        # 일반적인 횟수 패턴 감지 (예: "1-2회", "주 3-4회")
        import re
        if re.match(r'^[0-9]+[-~][0-9]+회$', text) or re.match(r'^주 [0-9]+[-~][0-9]+회$', text):
            return True
        
        # 소수점 및 숫자만 있는지 최종 확인
        try:
            float(text.replace(',', ''))
            return True
        except ValueError:
            return False
    
    def translate_ppt(self, ppt_path, source_lang, target_lang, vision_model, text_model, 
                    progress_callback=None, status_callback=None, debug_mode=False):
        """파워포인트 파일 번역 실행 (디버그 모드 추가)"""
        if debug_mode:
            # 디버그 모드에서는 로깅 레벨 변경
            original_level = logger.level
            logger.setLevel(logging.DEBUG)
            logger.info("디버그 모드 활성화됨")
        
        try:
            logger.info(f"번역 프로세스 시작: {ppt_path}")
            
            if status_callback:
                status_callback("번역 프로세스 시작")
            
            # 모델 상태 확인 추가
            if status_callback:
                status_callback("모델 상태 확인 중...")
            model_loaded, _ = self.ollama_service.check_ollama_model_status(vision_model)
            
            if not model_loaded:
                if status_callback:
                    status_callback(f"{vision_model} 모델 로드 중...")
                self.ollama_service.warmup_vision_model(vision_model)
            
            # 파워포인트 파일 열기
            ppt = Presentation(ppt_path)
            
            # 임시 폴더 생성
            with tempfile.TemporaryDirectory() as temp_dir:
                logger.info(f"임시 폴더 생성: {temp_dir}")
                
                # 문서 분석 (이미 분석된 경우에는 생략 가능)
                from services.document_analyzer import DocumentAnalyzer
                analyzer = DocumentAnalyzer()
                result = analyzer.analyze_ppt(ppt_path)
                
                text_elements = result['text_elements']
                image_elements = result['image_elements']
                total_elements = result['total_elements']
                
                # 프로그레스 초기화
                processed_items = 0
                
                # 1. 텍스트 요소 번역
                if status_callback:
                    status_callback("텍스트 요소 번역 중...")
                logger.info("텍스트 요소 번역 시작")
                
                for idx, text_element in enumerate(text_elements):
                    slide_idx = text_element['slide_idx']
                    slide = ppt.slides[slide_idx]
                    
                    try:
                        # 텍스트 요소 처리 부분
                        if text_element['type'] == 'text_run':
                            # 기존 로그 유지
                            logger.info(f"텍스트 번역 (슬라이드 {slide_idx+1}, 요소 {text_element['shape_idx']}): '{text_element['text'][:30]}...'")
                            
                            # 해당 텍스트 요소 찾기
                            shape = slide.shapes[text_element['shape_idx']]
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    # 수정: cell.text -> run.text
                                    if run.text.strip() == text_element['text']:
                                        # 숫자만 있는 텍스트인지 확인 (cell.text -> run.text)
                                        if self.is_numeric_text(run.text):
                                            logger.info(f"숫자 텍스트 감지됨, 번역 건너뜀: '{run.text}'")
                                            text_element['translated'] = True
                                        else:
                                            # 번역
                                            translated_text = self.ollama_service.translate_text(
                                                run.text, source_lang, target_lang, text_model
                                            )
                                            run.text = translated_text
                                            text_element['translated'] = True
                                        break
                        
                        elif text_element['type'] == 'table_cell':
                            # 기존 로그 유지
                            logger.info(f"테이블 셀 번역 (슬라이드 {slide_idx+1}, 테이블 {text_element['shape_idx']}, 행 {text_element['row_idx']}, 열 {text_element['col_idx']}): '{text_element['text'][:30]}...'")
                            
                            # 해당 테이블 찾기
                            shape = slide.shapes[text_element['shape_idx']]
                            if hasattr(shape, "table"):
                                table = shape.table
                                cell = table.rows[text_element['row_idx']].cells[text_element['col_idx']]
                                
                                # 번역
                                if cell.text.strip() == text_element['text']:
                                    # 숫자만 있는 텍스트인지 확인
                                    if self.is_numeric_text(cell.text):
                                        logger.info(f"숫자 텍스트 감지됨, 번역 건너뜀: '{cell.text}'")
                                        text_element['translated'] = True
                                    else:
                                        translated_text = self.ollama_service.translate_text(
                                            cell.text, source_lang, target_lang, text_model
                                        )
                                        
                                        # 텍스트 설정 (레퍼런스 이슈 방지)
                                        text_frame = cell.text_frame
                                        # 모든 기존 단락 제거
                                        while len(text_frame.paragraphs) > 1:
                                            tr_element = text_frame._txBody.remove(text_frame._txBody[1])
                                        
                                        # 첫 번째 단락에 새 텍스트 설정
                                        text_frame.paragraphs[0].text = translated_text
                                        
                                        text_element['translated'] = True
                    
                    except Exception as e:
                        logger.error(f"텍스트 번역 오류 (요소 {idx+1}/{len(text_elements)}): {str(e)}")
                        logger.error(traceback.format_exc())
                    
                    # 진행 상황 업데이트
                    processed_items += 1
                    if progress_callback:
                        progress_callback(processed_items, total_elements)
                
                # 2. 이미지 요소 번역
                if status_callback:
                    status_callback("이미지 요소 번역 중...")
                logger.info("이미지 요소 번역 시작")
                
                for idx, image_element in enumerate(image_elements):
                    slide_idx = image_element['slide_idx']
                    slide = ppt.slides[slide_idx]
                    
                    try:
                        logger.info(f"이미지 번역 (슬라이드 {slide_idx+1}, 요소 {image_element['shape_idx']})")
                        
                        # 해당 이미지 찾기
                        shape = slide.shapes[image_element['shape_idx']]
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # 이미지 추출 및 저장 시간 측정
                            save_start_time = time.time()
                            
                            # 이미지 추출 및 저장
                            image = shape.image
                            image_bytes = image.blob
                            
                            temp_image_path = os.path.join(temp_dir, f"slide_{slide_idx}_image_{image_element['shape_idx']}.png")
                            with open(temp_image_path, "wb") as f:
                                f.write(image_bytes)
                            
                            save_time = time.time() - save_start_time
                            logger.debug(f"이미지 저장 시간: {save_time:.3f}초")
                            logger.info(f"이미지 저장: {temp_image_path} ({len(image_bytes)} 바이트)")
                            
                            # 이미지 크기 확인 및 필요 시 리사이징 시간 측정
                            resize_start_time = time.time()
                            temp_image_path = resize_image_if_needed(temp_image_path)
                            resize_time = time.time() - resize_start_time
                            logger.debug(f"이미지 리사이징 시간: {resize_time:.3f}초")
                            
                            # 이미지를 Base64로 인코딩 시간 측정
                            encode_start_time = time.time()
                            image_base64 = encode_image_to_base64(temp_image_path)  # 이미 import한 함수 사용
                            encode_time = time.time() - encode_start_time
                            logger.debug(f"이미지 인코딩 시간: {encode_time:.3f}초")
                            logger.debug(f"인코딩된 이미지 크기: {len(image_base64)} 문자")
                            
                            # Vision 모델로 텍스트 추출 시간 측정
                            extract_start_time = time.time()
                            logger.info("Vision API 호출 시작")
                            extracted_text = self.ollama_service.extract_text_from_image(image_base64, vision_model)
                            extract_time = time.time() - extract_start_time
                            logger.info(f"Vision API 호출 완료: {extract_time:.3f}초")
                            
                            if extracted_text:
                                # 추출된 텍스트 번역
                                translated_text = self.ollama_service.translate_text(
                                    extracted_text, source_lang, target_lang, text_model
                                )
                                
                                if translated_text and translated_text != extracted_text:
                                    # 번역된 텍스트로 이미지 오버레이
                                    translated_image_path = overlay_text_on_image(temp_image_path, translated_text)
                                    
                                    # 번역된 이미지로 교체
                                    try:
                                        if os.path.exists(translated_image_path) and translated_image_path != temp_image_path:
                                            # 이미지 교체 방법 (Python-pptx 제한 우회)
                                            # 1. 새 이미지 슬라이드에 추가
                                            left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                            pic = slide.shapes.add_picture(translated_image_path, left, top, width, height)
                                            
                                            # 2. 기존 이미지와 동일한 위치/크기 설정
                                            pic.left, pic.top, pic.width, pic.height = left, top, width, height
                                            
                                            image_element['translated'] = True
                                            logger.info("이미지 교체 완료")
                                    except Exception as e:
                                        logger.error(f"이미지 교체 오류: {str(e)}")
                    
                    except Exception as e:
                        logger.error(f"이미지 번역 오류 (요소 {idx+1}/{len(image_elements)}): {str(e)}")
                        logger.error(traceback.format_exc())
                    
                    # 진행 상황 업데이트
                    processed_items += 1
                    if progress_callback:
                        progress_callback(processed_items, total_elements)
                
                # 번역된 파일 저장
                output_path = os.path.splitext(ppt_path)[0] + "_translated.pptx"
                logger.info(f"번역된 파일 저장: {output_path}")
                ppt.save(output_path)
                
                if status_callback:
                    status_callback(f"번역 완료! 파일 저장됨: {output_path}")
                
                logger.info(f"번역 완료")
                return output_path
            
        except Exception as e:
            logger.exception(f"번역 프로세스 오류: {str(e)}")
            if status_callback:
                status_callback(f"번역 오류: {str(e)}")
            raise
        finally:
            # 디버그 모드였다면 로깅 레벨 복원
            if debug_mode:
                logger.setLevel(original_level)
                logger.info("디버그 모드 비활성화됨")
