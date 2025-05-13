# services/translation.py
import os
import tempfile
import time
import logging
import traceback
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from utils.image_utils import resize_image_if_needed, overlay_text_on_image
from utils.image_utils import is_numeric_text

logger = logging.getLogger(__name__)

class TranslationService:
    def __init__(self, ollama_service):
        self.ollama_service = ollama_service

    def translate_ppt(self, ppt_path, source_lang, target_lang, text_model, 
                    progress_callback=None, status_callback=None, options=None):
        """파워포인트 파일 번역 실행"""
        if options is None:
            options = {}
        
        debug_mode = options.get('debug_mode', False)
        
        # 임시 파일 추적 리스트
        temp_files = []
        
        if debug_mode:
            original_level = logger.level
            logger.setLevel(logging.DEBUG)
            logger.info("디버그 모드 활성화됨")
        
        try:
            logger.info(f"번역 프로세스 시작: {ppt_path}")
            
            if status_callback:
                status_callback("번역 프로세스 시작")
            
            logger.info(f"사용할 모델: {text_model}")
            if status_callback:
                status_callback(f"번역 준비 중: {text_model}")
            
            # 파워포인트 파일 열기
            ppt = Presentation(ppt_path)
            
            # 임시 폴더 생성
            with tempfile.TemporaryDirectory() as temp_dir:
                logger.info(f"임시 폴더 생성: {temp_dir}")
                
                # 문서 분석
                from services.document_analyzer import DocumentAnalyzer
                analyzer = DocumentAnalyzer()
                result = analyzer.analyze_ppt(ppt_path)
                
                text_elements = result['text_elements']
                image_elements = result['image_elements']
                total_elements = result['total_elements']
                
                # 번역 작업 시작
                processed_items = 0
                
                # 1. 텍스트 요소 번역
                if status_callback:
                    status_callback("텍스트 요소 번역 중...")
                logger.info("텍스트 요소 번역 시작")
                
                self._translate_text_elements(
                    ppt, text_elements, source_lang, target_lang, text_model, 
                    progress_callback, processed_items, total_elements
                )
                processed_items += len(text_elements)
                
                # 2. 이미지 요소 번역
                if status_callback:
                    status_callback("이미지 요소 번역 중...")
                logger.info("이미지 요소 번역 시작")
                
                # 이미지 번역 로직에 temp_files 리스트 전달
                self._translate_image_elements(
                    ppt, image_elements, temp_dir, source_lang, target_lang, 
                    text_model, progress_callback, processed_items, total_elements, 
                    options, temp_files  # temp_files 추가
                )
                
                # 번역된 파일 저장
                output_path = os.path.splitext(ppt_path)[0] + "_translated.pptx"
                logger.info(f"번역된 파일 저장: {output_path}")
                ppt.save(output_path)
                
                # 임시 이미지 파일 삭제
                self._cleanup_temp_files(temp_files)
                
                if status_callback:
                    status_callback(f"번역 완료! 파일 저장됨: {output_path}")
                
                logger.info(f"번역 완료")
                return output_path
            
        except Exception as e:
            logger.exception(f"번역 프로세스 오류: {str(e)}")
            if status_callback:
                status_callback(f"번역 오류: {str(e)}")
            raise
        finally:
            # 실패 시에도 임시 파일 정리 시도
            if 'temp_files' in locals() and temp_files:
                self._cleanup_temp_files(temp_files)
                
            if debug_mode:
                logger.setLevel(original_level)
                logger.info("디버그 모드 비활성화됨")
    
    def _translate_text_elements(self, ppt, text_elements, source_lang, target_lang, text_model, 
                               progress_callback=None, processed_items=0, total_elements=0):
        """텍스트 요소 번역 처리"""
        for idx, text_element in enumerate(text_elements):
            slide_idx = text_element['slide_idx']
            slide = ppt.slides[slide_idx]
            
            try:
                # 텍스트 요소 타입에 따라 처리
                element_type = text_element['type']
                
                if element_type == 'paragraph':
                    self._translate_paragraph(slide, text_element, source_lang, target_lang, text_model)
                elif element_type == 'table_cell':
                    self._translate_table_cell(slide, text_element, source_lang, target_lang, text_model)
                elif element_type == 'text_run':
                    self._translate_text_run(slide, text_element, source_lang, target_lang, text_model)
                
            except Exception as e:
                logger.error(f"텍스트 번역 오류 (요소 {idx+1}/{len(text_elements)}): {str(e)}")
                logger.debug(traceback.format_exc())
            
            # 진행 상황 업데이트
            if progress_callback:
                current = processed_items + idx + 1
                progress_callback(current, total_elements)
    
    def _translate_paragraph(self, slide, text_element, source_lang, target_lang, text_model):
        """문단 번역 처리"""
        shape = slide.shapes[text_element['shape_idx']]
        paragraph = shape.text_frame.paragraphs[text_element['para_idx']]
        
        if paragraph.text.strip() == text_element['text']:
            if is_numeric_text(paragraph.text):
                logger.info(f"숫자 텍스트 감지됨, 번역 건너뜀: '{paragraph.text}'")
                text_element['translated'] = True
            else:
                translated_text = self.ollama_service.translate_text(
                    paragraph.text, source_lang, target_lang, text_model
                )
                
                # 서식 유지를 위한 처리
                if len(paragraph.runs) > 0:
                    # 첫 번째 run에 번역된 텍스트 설정
                    first_run = paragraph.runs[0]
                    first_run.text = translated_text
                    
                    # 나머지 run 제거
                    while len(paragraph.runs) > 1:
                        paragraph._p.remove(paragraph.runs[1]._r)
                else:
                    paragraph.text = translated_text
                
                text_element['translated'] = True
    
    def _translate_table_cell(self, slide, text_element, source_lang, target_lang, text_model):
        """테이블 셀 번역 처리"""
        shape = slide.shapes[text_element['shape_idx']]
        if hasattr(shape, "table"):
            table = shape.table
            cell = table.rows[text_element['row_idx']].cells[text_element['col_idx']]
            
            if cell.text.strip() == text_element['text']:
                if is_numeric_text(cell.text):
                    logger.info(f"숫자 텍스트 감지됨, 번역 건너뜀: '{cell.text}'")
                    text_element['translated'] = True
                else:
                    translated_text = self.ollama_service.translate_text(
                        cell.text, source_lang, target_lang, text_model
                    )
                    
                    # 서식 유지를 위한 처리
                    text_frame = cell.text_frame
                    if text_frame.paragraphs and len(text_frame.paragraphs[0].runs) > 0:
                        # 첫 번째 paragraph의 첫 번째 run에 번역된 텍스트 설정
                        first_run = text_frame.paragraphs[0].runs[0]
                        first_run.text = translated_text
                        
                        # 첫 번째 paragraph의 나머지 run 제거
                        while len(text_frame.paragraphs[0].runs) > 1:
                            text_frame.paragraphs[0]._p.remove(text_frame.paragraphs[0].runs[1]._r)
                        
                        # 첫 번째 이외의 paragraph 제거
                        while len(text_frame.paragraphs) > 1:
                            text_frame._txBody.remove(text_frame._txBody[1])
                    else:
                        # run이 없는 경우 직접 텍스트 설정
                        if text_frame.paragraphs:
                            text_frame.paragraphs[0].text = translated_text
                    
                    text_element['translated'] = True
    
    def _translate_text_run(self, slide, text_element, source_lang, target_lang, text_model):
        """텍스트 런 번역 처리 (레거시 지원)"""
        shape = slide.shapes[text_element['shape_idx']]
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() == text_element['text']:
                    if is_numeric_text(run.text):
                        logger.info(f"숫자 텍스트 감지됨, 번역 건너뜀: '{run.text}'")
                        text_element['translated'] = True
                    else:
                        translated_text = self.ollama_service.translate_text(
                            run.text, source_lang, target_lang, text_model
                        )
                        run.text = translated_text
                        text_element['translated'] = True
                    break
    
    def _translate_image_elements(self, ppt, image_elements, temp_dir, source_lang, target_lang,
                                text_model, progress_callback=None, processed_items=0, total_elements=0,
                                options=None, temp_files=None):
        """이미지 요소 번역 처리"""
        if options is None:
            options = {}
            
        if temp_files is None:
            temp_files = []
            
        debug_mode = options.get('debug_mode', False)
        source_lang_for_ocr = options.get('source_lang', source_lang)
        
        for idx, image_element in enumerate(image_elements):
            slide_idx = image_element['slide_idx']
            slide = ppt.slides[slide_idx]
            
            try:
                logger.info(f"이미지 번역 (슬라이드 {slide_idx+1}, 요소 {image_element['shape_idx']})")
                
                shape = slide.shapes[image_element['shape_idx']]
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 이미지 추출 및 임시 저장
                    image = shape.image
                    image_bytes = image.blob
                    
                    # 이미지 크기 확인 (너무 큰 이미지 건너뛰기)
                    if len(image_bytes) > 5*1024*1024:  # 5MB 제한
                        logger.warning(f"이미지 크기가 너무 큽니다 ({len(image_bytes)} 바이트). 건너뜁니다.")
                        continue
                    
                    # 임시 이미지 파일 저장
                    temp_image_path = os.path.join(temp_dir, f"slide_{slide_idx}_image_{image_element['shape_idx']}.png")
                    with open(temp_image_path, "wb") as f:
                        f.write(image_bytes)
                    
                    # 임시 파일 추적 목록에 추가
                    temp_files.append(temp_image_path)
                    logger.info(f"이미지 저장: {temp_image_path} ({len(image_bytes)} 바이트)")
                    
                    # 이미지 처리: 리사이징 및 OCR
                    temp_image_path = resize_image_if_needed(temp_image_path)
                    
                    try:
                        import pytesseract
                        import cv2
                        
                        # OCR 실행
                        img = cv2.imread(temp_image_path)
                        if img is None:
                            logger.error(f"이미지 로드 실패: {temp_image_path}")
                            continue
                            
                        # OCR 언어 설정 (config.py의 OCR_LANG_MAPPING 참조)
                        from config import OCR_LANG_MAPPING
                        
                        default_lang = 'eng'
                        ocr_lang = default_lang
                        
                        # 소스 언어에 따른 OCR 언어 설정
                        if source_lang_for_ocr in OCR_LANG_MAPPING:
                            ocr_lang = '+'.join(OCR_LANG_MAPPING[source_lang_for_ocr])
                        
                        logger.info(f"OCR 언어 설정: {ocr_lang}")
                        
                        # OCR 실행 (여러 PSM 모드 시도)
                        extracted_text = None
                        
                        for psm in [11, 6, 3, 4]:  # 다양한 PSM 모드 시도
                            config = f'--oem 3 --psm {psm}'
                            try:
                                text = pytesseract.image_to_string(img, lang=ocr_lang, config=config)
                                if text and text.strip():
                                    extracted_text = text
                                    logger.info(f"OCR 성공 (PSM {psm}): {len(text.strip())}자 추출")
                                    break
                            except Exception as e:
                                logger.error(f"OCR 오류 (PSM {psm}): {e}")
                                continue
                        
                        # 텍스트 추출 결과 확인
                        if not extracted_text or not extracted_text.strip():
                            logger.warning("OCR 실패: 텍스트를 추출할 수 없습니다.")
                            continue
                            
                        logger.info(f"추출된 텍스트: '{extracted_text.strip()}'")
                        
                        # 텍스트 번역
                        translated_text = self.ollama_service.translate_text(
                            extracted_text, source_lang, target_lang, text_model
                        )
                        
                        logger.info(f"번역된 텍스트: '{translated_text[:100]}'")
                        
                        if translated_text and translated_text != extracted_text:
                            # 번역된 텍스트로 이미지 오버레이
                            timestamp = int(time.time() * 1000)
                            translated_image_path = overlay_text_on_image(
                                temp_image_path, 
                                translated_text,
                                source_lang
                            )
                            
                            # 번역된 이미지 파일 추적
                            if translated_image_path != temp_image_path:
                                temp_files.append(translated_image_path)
                            
                            # 번역된 이미지로 교체
                            if os.path.exists(translated_image_path) and translated_image_path != temp_image_path:
                                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                try:
                                    # 기존 이미지 대신 새 이미지 추가
                                    pic = slide.shapes.add_picture(translated_image_path, left, top, width, height)
                                    # 위치와 크기 조정
                                    pic.left, pic.top, pic.width, pic.height = left, top, width, height
                                    
                                    image_element['translated'] = True
                                    logger.info("이미지 교체 완료")
                                except Exception as e:
                                    logger.error(f"이미지 교체 오류: {e}")
                        else:
                            logger.warning("텍스트가 번역되지 않았거나 원본과 동일합니다.")
                    
                    except ImportError:
                        logger.error("pytesseract 또는 OpenCV 모듈이 설치되지 않았습니다.")
                    
                    except Exception as e:
                        logger.error(f"이미지 OCR 처리 오류: {e}")
                        logger.debug(traceback.format_exc())
                    
                    # 원본 임시 파일 삭제 (즉시 삭제 옵션)
                    try:
                        if os.path.exists(temp_image_path):
                            os.remove(temp_image_path)
                            if temp_image_path in temp_files:
                                temp_files.remove(temp_image_path)  # 목록에서 제거
                    except Exception as e:
                        logger.debug(f"임시 파일 삭제 실패: {e}")
            
            except Exception as e:
                logger.error(f"이미지 번역 오류 (요소 {idx+1}/{len(image_elements)}): {str(e)}")
                logger.debug(traceback.format_exc())
            
            # 진행 상황 업데이트
            if progress_callback:
                current = processed_items + idx + 1
                progress_callback(current, total_elements)

    def _cleanup_temp_files(self, file_list):
        """임시 파일 정리"""
        count = 0
        for file_path in file_list:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    count += 1
            except Exception as e:
                logger.warning(f"임시 파일 삭제 실패: {file_path}, 오류: {e}")
        
        logger.info(f"{count}개의 임시 이미지 파일 삭제됨")